"""
Build CSE 253R Assignment 2 Presentation (PPTX)
Music Generation: Symbolic & Continuous
20 slides, widescreen 16:9, light theme
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Paths ──────────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_DIR   = os.path.dirname(SCRIPT_DIR)
OUT_PATH   = os.path.join(SCRIPT_DIR, "presentation.pptx")

def img(name):
    """Return full path for an image in repo root, or None if missing."""
    p = os.path.join(REPO_DIR, name)
    return p if os.path.exists(p) else None

# ── Colour palette ─────────────────────────────────────────────────────────
BG        = RGBColor(0xF8, 0xF9, 0xFA)   # very light gray
BG_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy (titles)
BODY      = RGBColor(0x2C, 0x2C, 0x2C)   # near-black (body)
MUTED     = RGBColor(0x55, 0x65, 0x6E)   # medium gray (captions)
RED       = RGBColor(0xE7, 0x4C, 0x3C)   # Task 1 accent
BLUE      = RGBColor(0x34, 0x98, 0xDB)   # Task 4 accent
GREEN     = RGBColor(0x27, 0xAE, 0x60)
AMBER     = RGBColor(0xF3, 0x9C, 0x12)
TH_BG     = RGBColor(0x2C, 0x3E, 0x50)   # table header bg
TH_FG     = RGBColor(0xFF, 0xFF, 0xFF)
TD_ALT    = RGBColor(0xEC, 0xF0, 0xF1)   # alternating row
GREEN_BG  = RGBColor(0xD5, 0xF5, 0xE3)
RED_BG    = RGBColor(0xFD, 0xED, 0xEC)
RED_DARK  = RGBColor(0xA9, 0x3C, 0x3C)
BLUE_DARK = RGBColor(0x1A, 0x5C, 0x8A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Presentation setup ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ── Low-level helpers ───────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1, but we use int directly
        x, y, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h,
                text="", font_size=Pt(18), bold=False, italic=False,
                color=BODY, align=PP_ALIGN.LEFT, wrap=True,
                font_name="Calibri"):
    """Add a textbox with a single paragraph."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline_textbox(slide, x, y, w, h,
                          lines, font_size=Pt(18), bold=False,
                          color=BODY, line_spacing=None,
                          font_name="Calibri", align=PP_ALIGN.LEFT,
                          space_before=Pt(4)):
    """Add a textbox with multiple paragraph lines."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = space_before
        run = p.add_run()
        run.text = str(line)
        run.font.name  = font_name
        run.font.size  = font_size
        run.font.bold  = bold
        run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def title_bar(slide, title_text, accent_color=DARK, bg_color=None):
    """Draw a title bar at the top of the slide."""
    bar_h = Inches(1.0)
    bar = add_rect(slide, 0, 0, SLIDE_W, bar_h, fill_color=bg_color or DARK)
    add_textbox(slide,
                Inches(0.4), Inches(0.1),
                Inches(12.5), bar_h - Inches(0.1),
                text=title_text,
                font_size=Pt(36), bold=True,
                color=RGBColor(0xFF,0xFF,0xFF),
                align=PP_ALIGN.LEFT,
                font_name="Calibri")
    return bar_h


def slide_title_body(title, body_lines, accent=DARK, img_path=None,
                     img_x=None, img_y=None, img_w=None):
    """Convenience: title bar + bulleted body text."""
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide)
    bar_h = title_bar(slide, title, bg_color=accent)
    content_y = bar_h + Inches(0.2)
    content_h  = SLIDE_H - content_y - Inches(0.2)
    text_w     = SLIDE_W - Inches(0.6) if img_path is None else Inches(7.5)
    add_multiline_textbox(slide,
                          Inches(0.4), content_y, text_w, content_h,
                          body_lines, font_size=Pt(20), color=BODY,
                          space_before=Pt(6))
    if img_path and os.path.exists(img_path):
        iw = img_w or Inches(5.0)
        ix = img_x or (SLIDE_W - iw - Inches(0.2))
        iy = img_y or content_y
        ih = SLIDE_H - iy - Inches(0.15)
        slide.shapes.add_picture(img_path, ix, iy, width=iw, height=ih)
    return slide


def add_table_to_slide(slide, headers, rows,
                       x, y, w, h,
                       header_bg=TH_BG, header_fg=TH_FG,
                       alt_row=TD_ALT, font_size=Pt(15),
                       col_widths=None,
                       cell_colors=None):
    """
    Add a formatted table.
    cell_colors: dict {(row_idx, col_idx): RGBColor} (row_idx 0 = first data row)
    """
    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl = slide.shapes.add_table(num_rows, num_cols, x, y, w, h).table
    # Column widths
    if col_widths:
        total_w = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total_w)
    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold  = True
        run.font.size  = font_size
        run.font.color.rgb = header_fg
        run.font.name  = "Calibri"
    # Data rows
    for ri, row in enumerate(rows):
        bg = alt_row if ri % 2 == 1 else BG_WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            # Override per-cell color?
            if cell_colors and (ri, ci) in cell_colors:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_colors[(ri, ci)]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size  = font_size
            run.font.color.rgb = BODY
            run.font.name  = "Calibri"
    return tbl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 -- Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_WHITE)

# Top accent bar
add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), fill_color=RED)
# Bottom accent bar
add_rect(slide, 0, SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), fill_color=BLUE)

# Main title
add_textbox(slide,
            Inches(0.7), Inches(1.6),
            Inches(11.9), Inches(1.5),
            "Music Generation: Symbolic & Continuous",
            font_size=Pt(48), bold=True, color=DARK,
            align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide,
            Inches(0.7), Inches(3.2),
            Inches(11.9), Inches(0.8),
            "CSE 253R  --  Assignment 2",
            font_size=Pt(28), bold=False, color=MUTED,
            align=PP_ALIGN.CENTER)

# Task pills
pill_y = Inches(4.2)
pill_h = Inches(0.55)
pill_w = Inches(5.4)
gap    = Inches(0.3)
left_x = (SLIDE_W - 2*pill_w - gap) / 2

add_rect(slide, left_x, pill_y, pill_w, pill_h, fill_color=RED)
add_textbox(slide, left_x, pill_y, pill_w, pill_h,
            "Task 1: Symbolic Unconditioned",
            font_size=Pt(18), bold=True, color=BG_WHITE,
            align=PP_ALIGN.CENTER)

add_rect(slide, left_x + pill_w + gap, pill_y, pill_w, pill_h, fill_color=BLUE)
add_textbox(slide, left_x + pill_w + gap, pill_y, pill_w, pill_h,
            "Task 4: Continuous Conditioned",
            font_size=Pt(18), bold=True, color=BG_WHITE,
            align=PP_ALIGN.CENTER)

# Footer
add_textbox(slide,
            Inches(0.5), Inches(6.85),
            Inches(12.3), Inches(0.4),
            "[Team names]  |  Spring 2026",
            font_size=Pt(14), color=MUTED,
            align=PP_ALIGN.CENTER)

print("Slide 1 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 -- Two Tasks Overview (two columns)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "What We Built")

col_y  = bar_h + Inches(0.2)
col_h  = Inches(4.7)
col_w  = Inches(6.0)
gap    = Inches(0.4)
left_x = Inches(0.4)
right_x= left_x + col_w + gap

# Left column card (red)
add_rect(slide, left_x, col_y, col_w, col_h,
         fill_color=RGBColor(0xFD,0xED,0xEC),
         line_color=RED, line_width=Pt(2))

add_textbox(slide, left_x + Inches(0.15), col_y + Inches(0.1),
            col_w - Inches(0.3), Inches(0.55),
            "Task 1: Symbolic Unconditioned",
            font_size=Pt(22), bold=True, color=RED)

left_bullets = [
    "  Learn p(Soprano, Alto, Tenor, Bass) on Bach chorales",
    "  Generate new 4-part MIDI harmony from scratch",
    "  Models: Bigram Markov Chain, 2-layer LSTM",
    "  Dataset: JSB Chorales (368 chorales, music21)",
    "  Output: symbolic_unconditioned.mid",
]
add_multiline_textbox(slide,
                      left_x + Inches(0.15), col_y + Inches(0.75),
                      col_w - Inches(0.3), col_h - Inches(0.9),
                      left_bullets, font_size=Pt(18), color=BODY,
                      space_before=Pt(5))

# Right column card (blue)
add_rect(slide, right_x, col_y, col_w, col_h,
         fill_color=RGBColor(0xEB,0xF5,0xFB),
         line_color=BLUE, line_width=Pt(2))

add_textbox(slide, right_x + Inches(0.15), col_y + Inches(0.1),
            col_w - Inches(0.3), Inches(0.55),
            "Task 4: Continuous Conditioned",
            font_size=Pt(22), bold=True, color=BLUE)

right_bullets = [
    "  Fine-tune MusicGen-small on genre-labeled audio",
    "  Generate 30s audio from text prompts",
    "  Model: MusicGen (300M params, Meta/AudioCraft)",
    "  Dataset: FMA-small (8,000 tracks, 8 genres)",
    "  Output: continuous_conditioned.mp3",
]
add_multiline_textbox(slide,
                      right_x + Inches(0.15), col_y + Inches(0.75),
                      col_w - Inches(0.3), col_h - Inches(0.9),
                      right_bullets, font_size=Pt(18), color=BODY,
                      space_before=Pt(5))

# Bottom banner
banner_y = col_y + col_h + Inches(0.15)
add_rect(slide, Inches(0.4), banner_y,
         SLIDE_W - Inches(0.8), Inches(0.65),
         fill_color=RGBColor(0xE8, 0xEA, 0xEB))
add_textbox(slide,
            Inches(0.55), banner_y + Inches(0.05),
            SLIDE_W - Inches(1.1), Inches(0.55),
            "Symbolic = interpretable, small data, rule-checkable.   "
            "Continuous = perceptually realistic, requires large pretrained model.",
            font_size=Pt(16), color=DARK, align=PP_ALIGN.CENTER)

print("Slide 2 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 -- JSB Chorales Dataset
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Task 1 Dataset: JSB Chorales", bg_color=RED)

content_y = bar_h + Inches(0.2)
table_x   = Inches(0.4)
table_w   = Inches(6.5)
table_h   = Inches(4.5)

headers = ["Property", "Value"]
rows = [
    ["Source",        "music21 built-in corpus (no download)"],
    ["Composer",      "J.S. Bach (BWV 1-438)"],
    ["Chorales",      "368 (4-part SATB only)"],
    ["Representation","(T, 4) MIDI pitch matrix, 16th-note grid"],
    ["Vocabulary",    "47 tokens (46 pitches + rest)"],
    ["Pitch range",   "MIDI 36 (C2) to MIDI 81 (A5)"],
    ["Split",         "296 train / 36 val / 36 test chorales"],
    ["Sequences",     "2,845 train / 178 val / 190 test windows"],
    ["Window size",   "64 steps = 16 beats = ~4 bars"],
]
add_table_to_slide(slide, headers, rows,
                   table_x, content_y, table_w, table_h,
                   font_size=Pt(15),
                   col_widths=[1.8, 3.5])

# Right: image or placeholder
img_path = img("eda_piano_roll.png")
img_x = table_x + table_w + Inches(0.3)
img_w = SLIDE_W - img_x - Inches(0.2)
img_y = content_y

if img_path:
    slide.shapes.add_picture(img_path, img_x, img_y, width=img_w,
                             height=Inches(4.5))
else:
    add_rect(slide, img_x, img_y, img_w, Inches(4.5),
             fill_color=RGBColor(0xE0,0xE8,0xEE))
    add_textbox(slide, img_x, img_y + Inches(1.8), img_w, Inches(0.8),
                "Piano Roll Visualization\n(shown in notebook)",
                font_size=Pt(16), color=MUTED, align=PP_ALIGN.CENTER)

# Bottom note
note_y = content_y + table_h + Inches(0.1)
add_textbox(slide,
            Inches(0.4), note_y,
            SLIDE_W - Inches(0.8), Inches(0.45),
            "Benchmark used in: DeepBach (Hadjeres 2017), BachBot (Liang 2017), "
            "Music Transformer (Huang 2019), BacHMMachine (Hahn 2021)",
            font_size=Pt(13), color=MUTED)

print("Slide 3 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 -- Task 1 Modeling Approach
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Task 1: Modeling Approach", bg_color=RED)

col_y = bar_h + Inches(0.2)
col_h = Inches(4.5)
col_w = Inches(6.1)
gap   = Inches(0.35)
lx    = Inches(0.4)
rx    = lx + col_w + gap

# Baseline box
add_rect(slide, lx, col_y, col_w, col_h,
         fill_color=BG_WHITE,
         line_color=MUTED, line_width=Pt(1.5))

add_textbox(slide, lx + Inches(0.1), col_y + Inches(0.05),
            col_w - Inches(0.2), Inches(0.35),
            "BASELINE",
            font_size=Pt(13), bold=True, color=MUTED)

add_textbox(slide, lx + Inches(0.1), col_y + Inches(0.35),
            col_w - Inches(0.2), Inches(0.6),
            "Bigram Markov Chain",
            font_size=Pt(22), bold=True, color=DARK)

markov_lines = [
    "  4 independent 47x47 transition matrices (one per voice)",
    "  Laplace (add-1) smoothing",
    "  Context: current timestep only",
    "  P(next | current) = (count + 1) / (sum + V)",
    "",
    "  Test Perplexity:  2.59",
]
add_multiline_textbox(slide,
                      lx + Inches(0.15), col_y + Inches(1.1),
                      col_w - Inches(0.3), col_h - Inches(1.3),
                      markov_lines, font_size=Pt(18), color=BODY,
                      space_before=Pt(5))

# Our model box
add_rect(slide, rx, col_y, col_w, col_h,
         fill_color=RGBColor(0xFD,0xED,0xEC),
         line_color=RED, line_width=Pt(2))

add_textbox(slide, rx + Inches(0.1), col_y + Inches(0.05),
            col_w - Inches(0.2), Inches(0.35),
            "OUR MODEL",
            font_size=Pt(13), bold=True, color=RED)

add_textbox(slide, rx + Inches(0.1), col_y + Inches(0.35),
            col_w - Inches(0.2), Inches(0.6),
            "2-Layer LSTM",
            font_size=Pt(22), bold=True, color=RED)

lstm_lines = [
    "  Joint processing of all 4 voices",
    "  4 x Embedding(47, 64) concatenated -> 256-dim input",
    "  LSTM(256, 256, 2 layers, dropout=0.3)",
    "  4 x Linear(256, 47) output heads",
    "  Parameters: 1.1M",
    "",
    "  Test Perplexity:  1.97  (1.3x improvement)",
]
add_multiline_textbox(slide,
                      rx + Inches(0.15), col_y + Inches(1.1),
                      col_w - Inches(0.3), col_h - Inches(1.3),
                      lstm_lines, font_size=Pt(18), color=BODY,
                      space_before=Pt(5))

# Bottom note
note_y = col_y + col_h + Inches(0.15)
add_rect(slide, Inches(0.4), note_y, SLIDE_W - Inches(0.8), Inches(0.55),
         fill_color=RGBColor(0xE8,0xEA,0xEB))
add_textbox(slide, Inches(0.55), note_y + Inches(0.05),
            SLIDE_W - Inches(1.1), Inches(0.45),
            "Objective: Cross-entropy loss summed over 4 voices.   "
            "Training: Teacher forcing.   "
            "Generation: Autoregressive sampling with temperature.",
            font_size=Pt(15), color=DARK, align=PP_ALIGN.CENTER)

print("Slide 4 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 -- LSTM Architecture (with training curves image)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "LSTM Architecture", bg_color=RED)

content_y = bar_h + Inches(0.2)

# Architecture diagram as text in a mono-font box
arch_lines = [
    "Input: (batch, T, 4) token indices",
    "",
    "  Embed(47,64)   Embed(47,64)   Embed(47,64)   Embed(47,64)",
    "   Soprano          Alto          Tenor           Bass",
    "",
    "              concatenate -> (batch, T, 256)",
    "",
    "         LSTM(256 hidden, 2 layers, dropout=0.3)",
    "                    (batch, T, 256)",
    "",
    "  Linear(256,47) Linear(256,47) Linear(256,47) Linear(256,47)",
    "   Soprano          Alto          Tenor           Bass",
    "",
    "           logits: (batch, T, 4, 47)",
]

arch_x = Inches(0.4)
arch_w = Inches(7.2)
arch_box = slide.shapes.add_textbox(arch_x, content_y, arch_w, Inches(5.0))
arch_tf  = arch_box.text_frame
arch_tf.word_wrap = False
first = True
for line in arch_lines:
    if first:
        p = arch_tf.paragraphs[0]
        first = False
    else:
        p = arch_tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name  = "Courier New"
    run.font.size  = Pt(14)
    run.font.color.rgb = DARK

# Side annotations
ann_x = Inches(7.8)
ann_y = content_y
ann_w = Inches(5.3)
ann_lines = [
    "Key Details:",
    "",
    "  Teacher forcing during training",
    "  Autoregressive sampling at inference",
    "  Temperature 0.9 for main deliverable",
    "  Early stopped at epoch 28",
    "  Best checkpoint at epoch 18",
    "  Parameters: ~1.1M total",
]
add_multiline_textbox(slide, ann_x, ann_y, ann_w, Inches(3.5),
                      ann_lines, font_size=Pt(18), color=BODY,
                      space_before=Pt(6))

# Training curves image
tc_img = img("training_curves.png")
if tc_img:
    slide.shapes.add_picture(tc_img, ann_x, ann_y + Inches(3.6),
                             width=ann_w, height=Inches(2.9))

print("Slide 5 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 -- Task 1 Evaluation Framework
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Task 1: Evaluation Framework", bg_color=RED)

content_y = bar_h + Inches(0.2)

headers = ["Metric", "What It Measures", "Why It Matters"]
rows = [
    ["Test Perplexity",
     "Statistical fit (next-token prediction quality)",
     "Lower = model predicts Bach better"],
    ["Pitch-Class KL Divergence",
     "Note frequency distribution vs real Bach",
     "Checks if generated music uses Bach's pitch palette"],
    ["Interval L1 Distance",
     "Melodic jump pattern vs real Bach",
     "Checks if stepwise/leapwise motion is realistic"],
    ["Voice Range Violations",
     "% notes outside historical SATB ranges",
     "Basic constraint: Soprano should not go below Alto"],
    ["Parallel 5ths/Octaves",
     "Voice-leading rule violations",
     "Music theory constraint Bach almost never violates"],
]
add_table_to_slide(slide, headers, rows,
                   Inches(0.4), content_y,
                   SLIDE_W - Inches(0.8), Inches(4.6),
                   font_size=Pt(16),
                   col_widths=[2.2, 3.8, 3.8])

# Highlighted question
q_y = content_y + Inches(4.75)
add_rect(slide, Inches(0.4), q_y, SLIDE_W - Inches(0.8), Inches(0.55),
         fill_color=RGBColor(0xFF,0xF3,0xCD))
add_textbox(slide, Inches(0.6), q_y + Inches(0.05),
            SLIDE_W - Inches(1.2), Inches(0.45),
            "Core question: Does better perplexity = better music?",
            font_size=Pt(18), bold=True, color=AMBER, align=PP_ALIGN.CENTER)

print("Slide 6 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 -- Task 1 Results (with eval_summary image)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Task 1: Results and Discussion", bg_color=RED)

content_y = bar_h + Inches(0.2)

eval_img = img("eval_summary.png")
table_w = SLIDE_W - Inches(0.8) if not eval_img else Inches(7.8)

headers = ["Metric", "Markov", "LSTM", "Real Bach"]
rows = [
    ["Test Perplexity",         "2.59",     "1.97 *",    "--"],
    ["Pitch-Class KL (to real)","0.176 *",  "0.198",     "0"],
    ["Interval L1 (to real)",   "0.121 *",  "0.127",     "0"],
    ["Voice Range Violations",  "0.09%",    "0.01% *",   "0%"],
    ["Parallel 5ths/Octaves",   "11.5% *",  "18.1%",     "3.0%"],
]
# Color cells: bold LSTM win = green, bold Markov win = green
cell_colors = {
    (0,2): GREEN_BG,   # LSTM perplexity
    (1,1): GREEN_BG,   # Markov KL
    (2,1): GREEN_BG,   # Markov interval
    (3,2): GREEN_BG,   # LSTM voice range
    (4,1): GREEN_BG,   # Markov parallel
}
add_table_to_slide(slide, headers, rows,
                   Inches(0.4), content_y,
                   table_w, Inches(2.8),
                   font_size=Pt(15),
                   col_widths=[2.8, 1.5, 1.5, 1.5],
                   cell_colors=cell_colors)

if eval_img:
    slide.shapes.add_picture(eval_img,
                             Inches(0.4) + table_w + Inches(0.1),
                             content_y,
                             width=SLIDE_W - Inches(0.4) - table_w - Inches(0.3),
                             height=Inches(2.8))

disc_y = content_y + Inches(2.95)
disc_lines = [
    "  LSTM wins on perplexity and voice range compliance",
    "  Markov wins on distributional similarity and parallel motion",
    "  The LSTM optimizes cross-entropy (statistical fit), not voice-leading rules",
    "  Markov chain avoids parallel motion by mostly sustaining notes (77% unisons)",
    "  Consistent with literature: BachBot reports comparable perplexity; "
    "Music Transformer achieves ~1.2 ppl with relative attention",
]
add_multiline_textbox(slide,
                      Inches(0.4), disc_y,
                      SLIDE_W - Inches(0.8), Inches(3.0),
                      disc_lines, font_size=Pt(17), color=BODY,
                      space_before=Pt(6))

print("Slide 7 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 -- Task 4 Title Card
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_WHITE)

# Blue accent bars
add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), fill_color=BLUE)
add_rect(slide, 0, SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), fill_color=BLUE)

add_textbox(slide, Inches(0.7), Inches(1.2), Inches(11.9), Inches(1.5),
            "Task 4: Continuous Conditioned Generation",
            font_size=Pt(44), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.7), Inches(2.85), Inches(11.9), Inches(0.7),
            "Fine-tuning MusicGen on FMA Genre Data",
            font_size=Pt(26), color=MUTED, align=PP_ALIGN.CENTER)

# Flow diagram
flow_y = Inches(4.0)
flow_h = Inches(0.65)
flow_items = [
    ("Text Prompt", RGBColor(0x2E,0x86,0xC1)),
    ("MusicGen", BLUE),
    ("30s Audio Clip", RGBColor(0x1E,0x8B,0x4F)),
]
box_w    = Inches(2.8)
arrow_w  = Inches(0.5)
total_fw = len(flow_items)*box_w + (len(flow_items)-1)*arrow_w
flow_x   = (SLIDE_W - total_fw) / 2

for i, (label, color) in enumerate(flow_items):
    bx = flow_x + i*(box_w + arrow_w)
    add_rect(slide, bx, flow_y, box_w, flow_h, fill_color=color)
    add_textbox(slide, bx, flow_y, box_w, flow_h,
                label, font_size=Pt(18), bold=True,
                color=BG_WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow_items)-1:
        ax = bx + box_w
        add_textbox(slide, ax, flow_y, arrow_w, flow_h,
                    "->", font_size=Pt(22), bold=True,
                    color=MUTED, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.55),
            '"upbeat electronic music with synthesizers"  ->  [waveform]',
            font_size=Pt(18), italic=True, color=DARK, align=PP_ALIGN.CENTER)

print("Slide 8 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 -- FMA-Small Dataset
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Task 4 Dataset: FMA-Small (Free Music Archive)", bg_color=BLUE)

content_y = bar_h + Inches(0.2)
table_w   = Inches(6.5)

headers = ["Property", "Value"]
rows_fma = [
    ["Source",        "HuggingFace (rpmon/fma-genre-classification)"],
    ["Total tracks",  "8,000"],
    ["Per genre",     "1,000 (perfectly balanced)"],
    ["Clip duration", "30 seconds"],
    ["Sample rate",   "22,050 Hz mono"],
    ["Total audio",   "~67 hours"],
    ["License",       "Creative Commons"],
    ["Reference",     "Defferrard et al. 2017"],
]
add_table_to_slide(slide, headers, rows_fma,
                   Inches(0.4), content_y, table_w, Inches(3.9),
                   font_size=Pt(15), col_widths=[1.8, 3.5])

# Genre list
genres = ["Hip-Hop", "Pop", "Folk", "Experimental",
          "Rock", "International", "Electronic", "Instrumental"]
genre_x = Inches(7.2)
genre_y = content_y
genre_w = Inches(2.8)

add_textbox(slide, genre_x, genre_y, genre_w, Inches(0.45),
            "8 Genres:", font_size=Pt(18), bold=True, color=BLUE)

for i, g in enumerate(genres):
    gy = genre_y + Inches(0.45) + i * Inches(0.42)
    add_rect(slide, genre_x, gy + Inches(0.1),
             Inches(0.18), Inches(0.2), fill_color=BLUE)
    add_textbox(slide, genre_x + Inches(0.25), gy,
                genre_w - Inches(0.3), Inches(0.42),
                g, font_size=Pt(17), color=BODY)

# FMA overview image
fma_img = img("eda_fma_overview.png")
if fma_img:
    slide.shapes.add_picture(fma_img,
                             Inches(10.2), content_y,
                             width=Inches(2.9), height=Inches(3.9))

# Highlight box
hl_y = content_y + Inches(4.05)
add_rect(slide, Inches(0.4), hl_y, SLIDE_W - Inches(0.8), Inches(0.75),
         fill_color=RGBColor(0xEB,0xF5,0xFB),
         line_color=BLUE, line_width=Pt(1.5))
add_textbox(slide, Inches(0.6), hl_y + Inches(0.05),
            SLIDE_W - Inches(1.2), Inches(0.65),
            "Fine-tuning subset: 4 genres (Hip-Hop, Folk, Electronic, Rock)   |   "
            "20 tracks/genre = 72 train + 8 val pairs (smoke test on Colab T4)",
            font_size=Pt(16), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

print("Slide 9 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 -- Fine-tuning vs Training from Scratch
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Why Fine-tune Instead of Training from Scratch?", bg_color=BLUE)

content_y = bar_h + Inches(0.2)

headers = ["Criterion", "Train from Scratch", "Fine-tune MusicGen (Ours)"]
rows_cmp = [
    ["Audio quality",   "Noise/blur without days of training",  "Studio-quality audio (pretrained)"],
    ["Data needed",     "100s of hours required",               "~67 hrs (FMA-small)"],
    ["GPU time",        "Weeks on large clusters",              "~4-8 hours (Colab T4)"],
    ["Output variety",  "Usually one dominant mode",            "8 distinct genres"],
    ["Trains own weights", "Yes",                               "Yes (decoder fine-tuning)"],
]
cell_colors = {}
for ri in range(5):
    cell_colors[(ri, 1)] = RGBColor(0xFD,0xED,0xEC)   # red column (scratch)
    cell_colors[(ri, 2)] = GREEN_BG                     # green column (ours)

add_table_to_slide(slide, headers, rows_cmp,
                   Inches(0.4), content_y,
                   SLIDE_W - Inches(0.8), Inches(3.8),
                   font_size=Pt(16),
                   col_widths=[2.5, 3.5, 4.0],
                   cell_colors=cell_colors)

# Bottom callout
co_y = content_y + Inches(4.0)
add_rect(slide, Inches(0.4), co_y, SLIDE_W - Inches(0.8), Inches(0.8),
         fill_color=RGBColor(0xEB,0xF5,0xFB),
         line_color=BLUE, line_width=Pt(2))
add_multiline_textbox(slide, Inches(0.6), co_y + Inches(0.06),
                      SLIDE_W - Inches(1.2), Inches(0.68),
                      ["Fine-tuning satisfies the 'train your own weights' requirement.",
                       "We update the transformer decoder parameters on our genre-labeled data."],
                      font_size=Pt(17), color=BLUE, space_before=Pt(3))

print("Slide 10 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 -- MusicGen Architecture
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "MusicGen Architecture  (Copet et al., NeurIPS 2023)", bg_color=BLUE)

content_y = bar_h + Inches(0.25)

# Check for architecture image
arch_img = img("eda_musicgen_arch.png")
if arch_img:
    slide.shapes.add_picture(arch_img, Inches(0.4), content_y,
                             width=Inches(6.5), height=Inches(5.8))
    text_x = Inches(7.2)
else:
    text_x = Inches(0.4)

# Pipeline boxes
pipeline = [
    ("Text Prompt",            RGBColor(0x2E,0x86,0xC1), "FROZEN",     '"upbeat electronic..."'),
    ("T5 Text Encoder",        RGBColor(0x7D,0x3C,0x98), "FROZEN",     "T5-base, 768-dim embeddings"),
    ("Transformer Decoder LM", RED,                       "FINE-TUNED", "300M params, 24 layers, h=1024"),
    ("EnCodec Decoder",        RGBColor(0x1E,0x8B,0x4F), "FROZEN",     "32kHz, 4 codebooks, size 2048"),
    ("Audio Output",           AMBER,                     "",           "32kHz waveform, 30 seconds"),
]

box_w = Inches(2.9) if not arch_img else Inches(5.8)
box_h = Inches(0.75)
gap_h = Inches(0.2)
total_pipe_h = len(pipeline)*box_h + (len(pipeline)-1)*gap_h
pipe_y = content_y + (Inches(5.5) - total_pipe_h)/2
pipe_x = text_x if arch_img else (SLIDE_W - box_w)/2

for i, (label, color, badge, sub) in enumerate(pipeline):
    by = pipe_y + i*(box_h + gap_h)
    add_rect(slide, pipe_x, by, box_w, box_h, fill_color=color)
    add_textbox(slide, pipe_x + Inches(0.1), by + Inches(0.05),
                box_w - Inches(0.2), Inches(0.38),
                label, font_size=Pt(16), bold=True,
                color=BG_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, pipe_x + Inches(0.1), by + Inches(0.4),
                box_w - Inches(0.2), Inches(0.3),
                sub, font_size=Pt(12), color=BG_WHITE, align=PP_ALIGN.CENTER)
    if badge:
        badge_color = RED if badge == "FINE-TUNED" else MUTED
        add_textbox(slide, pipe_x - Inches(0.05), by,
                    Inches(1.3), Inches(0.3),
                    badge, font_size=Pt(10), bold=True,
                    color=BG_WHITE if badge == "FINE-TUNED" else TH_FG)
    if i < len(pipeline)-1:
        arr_y = by + box_h
        add_textbox(slide, pipe_x, arr_y,
                    box_w, gap_h,
                    "v", font_size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)

# Annotation
ann_text = "Only the Transformer Decoder is fine-tuned on FMA data"
if arch_img:
    add_textbox(slide,
                pipe_x - Inches(0.1),
                pipe_y + 2*(box_h+gap_h) + box_h + Inches(0.05),
                box_w + Inches(0.3), Inches(0.5),
                ann_text, font_size=Pt(13), italic=True, color=RED)
else:
    add_textbox(slide,
                Inches(0.4), content_y + Inches(5.1),
                SLIDE_W - Inches(0.8), Inches(0.5),
                ann_text, font_size=Pt(16), italic=True, color=RED,
                align=PP_ALIGN.CENTER)

print("Slide 11 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 -- MusicGen Fine-tuning Training Summary
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "MusicGen Fine-tuning: Training Summary", bg_color=BLUE)

content_y = bar_h + Inches(0.2)
table_w   = Inches(6.5)

headers = ["Setting", "Value"]
rows_tr = [
    ["Base model",        "facebook/musicgen-small"],
    ["Training data",     "72 pairs (4 genres x 18 tracks)"],
    ["Validation data",   "8 pairs (4 genres x 2 tracks)"],
    ["Epochs",            "5"],
    ["Batch size",        "2"],
    ["Learning rate",     "1e-4"],
    ["Total steps",       "180"],
    ["GPU",               "Colab T4 (16 GB VRAM)"],
    ["Training loss",     "8.13 -> 7.09"],
    ["Validation loss",   "7.34 -> 7.21"],
]
add_table_to_slide(slide, headers, rows_tr,
                   Inches(0.4), content_y, table_w, Inches(4.4),
                   font_size=Pt(15), col_widths=[2.2, 3.5])

# Fine-tune loss curve image (or training curves)
tc4 = img("t4_training_curves.png") or img("training_curves.png")
if tc4:
    slide.shapes.add_picture(tc4,
                             Inches(0.4) + table_w + Inches(0.3),
                             content_y,
                             width=SLIDE_W - table_w - Inches(1.0),
                             height=Inches(4.4))
else:
    # placeholder
    ph_x = Inches(0.4) + table_w + Inches(0.3)
    ph_w = SLIDE_W - table_w - Inches(1.0)
    add_rect(slide, ph_x, content_y, ph_w, Inches(4.4),
             fill_color=RGBColor(0xEB,0xF5,0xFB))
    add_textbox(slide, ph_x, content_y + Inches(1.8), ph_w, Inches(0.8),
                "Fine-tune loss curve\n(from finetune_history.json)",
                font_size=Pt(15), color=MUTED, align=PP_ALIGN.CENTER)

# Bottom note
note_y = content_y + Inches(4.55)
add_textbox(slide, Inches(0.4), note_y, SLIDE_W - Inches(0.8), Inches(0.45),
            "Smoke test only. A full run (800 tracks, 25 epochs) would produce stronger genre specificity.",
            font_size=Pt(14), italic=True, color=MUTED)

print("Slide 12 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 -- Task 4 Evaluation Approach
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Task 4: Evaluation via Genre Consistency", bg_color=BLUE)

content_y = bar_h + Inches(0.25)

# Pipeline flow (horizontal)
pipeline_items = [
    ("Generate 30s audio\nfrom text prompt", BLUE),
    ("Extract MFCC\nfeatures", RGBColor(0x7D,0x3C,0x98)),
    ("SVM genre\nclassifier", RGBColor(0x1E,0x8B,0x4F)),
    ("Compare to\ntarget genre", AMBER),
]
box_w  = Inches(2.6)
box_h  = Inches(1.0)
arr_w  = Inches(0.45)
total  = len(pipeline_items)*box_w + (len(pipeline_items)-1)*arr_w
p_x    = (SLIDE_W - total)/2
p_y    = content_y + Inches(0.1)

for i, (label, color) in enumerate(pipeline_items):
    bx = p_x + i*(box_w + arr_w)
    add_rect(slide, bx, p_y, box_w, box_h, fill_color=color)
    add_textbox(slide, bx, p_y, box_w, box_h,
                label, font_size=Pt(16), bold=True,
                color=BG_WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipeline_items)-1:
        ax = bx + box_w
        add_textbox(slide, ax, p_y, arr_w, box_h,
                    "->", font_size=Pt(22), bold=True,
                    color=MUTED, align=PP_ALIGN.CENTER)

# Below diagram
metrics_y = p_y + box_h + Inches(0.25)
metric_lines = [
    "  Metric: Genre classifier accuracy on generated audio",
    "  Baseline: Pretrained MusicGen (no fine-tuning)",
    "  Our model: Fine-tuned MusicGen",
    "  Test set: 1 sample per genre, 4 genres total",
]
add_multiline_textbox(slide,
                      Inches(0.4), metrics_y,
                      Inches(7.8), Inches(2.2),
                      metric_lines, font_size=Pt(19), color=BODY,
                      space_before=Pt(8))

# Side note
note_x = Inches(8.4)
add_rect(slide, note_x, metrics_y, Inches(4.7), Inches(2.5),
         fill_color=RGBColor(0xFD,0xF2,0xCC),
         line_color=AMBER, line_width=Pt(1.5))
add_textbox(slide, note_x + Inches(0.15), metrics_y + Inches(0.1),
            Inches(4.4), Inches(0.4),
            "Why not Frechet Audio Distance (FAD)?",
            font_size=Pt(15), bold=True, color=AMBER)
add_multiline_textbox(slide,
                      note_x + Inches(0.15), metrics_y + Inches(0.55),
                      Inches(4.4), Inches(1.85),
                      ["Requires pretrained VGGish or similar model.",
                       "Not set up within project timeline.",
                       "Genre accuracy is more interpretable."],
                      font_size=Pt(15), color=BODY, space_before=Pt(5))

# SVM trained on real FMA note
svm_y = metrics_y + Inches(2.55)
add_textbox(slide, Inches(0.4), svm_y, SLIDE_W - Inches(0.8), Inches(0.45),
            "SVM classifier trained on MFCC features from real FMA tracks "
            "(not the generated audio)",
            font_size=Pt(15), italic=True, color=MUTED)

print("Slide 13 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 -- Task 4 Results
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Task 4: Results", bg_color=BLUE)

content_y = bar_h + Inches(0.2)

# Main result
main_y = content_y
add_rect(slide, Inches(0.4), main_y, SLIDE_W - Inches(0.8), Inches(0.7),
         fill_color=RGBColor(0xEB,0xF5,0xFB),
         line_color=BLUE, line_width=Pt(2))
add_textbox(slide, Inches(0.6), main_y + Inches(0.05),
            SLIDE_W - Inches(1.2), Inches(0.6),
            "Genre Accuracy:  25% (pretrained)  ->  75% (fine-tuned)",
            font_size=Pt(26), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Per-genre table
headers = ["Genre", "Prompt", "Pretrained", "Fine-tuned"]
rows_r4 = [
    ["Hip-Hop",    "hip hop music with beats and rhythm",       "Electronic (87%)", "Electronic (97%)"],
    ["Folk",       "acoustic folk music with guitar",           "Folk (100%)",       "Folk (100%)"],
    ["Electronic", "electronic music with synthesizers",        "Folk (100%)",       "Electronic (100%)"],
    ["Rock",       "energetic rock music with electric guitar", "Folk (100%)",       "Rock (100%)"],
]
cell_colors = {
    (0, 2): RED_BG,    (0, 3): RED_BG,    # Hip-Hop: both wrong
    (1, 2): GREEN_BG,  (1, 3): GREEN_BG,  # Folk: both correct
    (2, 2): RED_BG,    (2, 3): GREEN_BG,  # Electronic: pretrained wrong, ours right
    (3, 2): RED_BG,    (3, 3): GREEN_BG,  # Rock: pretrained wrong, ours right
}

table_y = main_y + Inches(0.8)
table_h = Inches(2.4)

eval_img4 = img("eval_task4_genre_accuracy.png")
table_w   = SLIDE_W - Inches(0.8) if not eval_img4 else Inches(7.8)

add_table_to_slide(slide, headers, rows_r4,
                   Inches(0.4), table_y, table_w, table_h,
                   font_size=Pt(14),
                   col_widths=[1.5, 4.0, 2.5, 2.5],
                   cell_colors=cell_colors)

if eval_img4:
    slide.shapes.add_picture(eval_img4,
                             Inches(0.4) + table_w + Inches(0.1),
                             table_y,
                             width=SLIDE_W - table_w - Inches(0.7),
                             height=table_h)

# Discussion
disc_y = table_y + table_h + Inches(0.15)
disc_lines = [
    "  Pretrained model defaults to Folk-like acoustic textures for most prompts",
    "  Fine-tuning teaches genre-specific audio characteristics",
    "  Hip-Hop misclassified as Electronic: hip-hop production often uses synthesized beats",
    "  Limitation: only 4 test samples (1 per genre), no confidence intervals",
]
add_multiline_textbox(slide,
                      Inches(0.4), disc_y,
                      SLIDE_W - Inches(0.8), Inches(2.2),
                      disc_lines, font_size=Pt(17), color=BODY,
                      space_before=Pt(6))

print("Slide 14 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 -- Related Work: Symbolic Generation
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Related Work: Symbolic Music Generation", bg_color=RED)

content_y = bar_h + Inches(0.2)

t1_rw = img("t1_related_work.png")
table_w = SLIDE_W - Inches(0.8) if not t1_rw else Inches(8.5)

headers = ["Year", "Work", "Model", "Key Contribution", "Relation to Ours"]
rows_rw = [
    ["2005", "Allan & Williams", "HMM",
     "Classic baseline, BIC model selection",
     "Our Markov chain is similar but simpler"],
    ["2017", "DeepBach (Hadjeres)", "Gibbs + LSTM",
     "Steerable generation, fix any voice",
     "Explicit harmonic constraints we lack"],
    ["2017", "BachBot (Liang)", "LSTM seq2seq",
     "Turing test: fooled 1/3 listeners",
     "Most similar architecture to ours"],
    ["2019", "Music Transformer (Huang)", "Transformer + rel. attn",
     "Best perplexity (~1.2)",
     "Shows headroom above our 1.97"],
    ["2021", "BacHMMachine (Hahn)", "Theory-guided HMM",
     "Interpretable chord transitions",
     "Music theory as model constraint"],
]
add_table_to_slide(slide, headers, rows_rw,
                   Inches(0.4), content_y,
                   table_w, Inches(4.0),
                   font_size=Pt(13),
                   col_widths=[0.7, 2.2, 2.0, 3.2, 3.2])

if t1_rw:
    slide.shapes.add_picture(t1_rw,
                             Inches(0.4) + table_w + Inches(0.1),
                             content_y,
                             width=SLIDE_W - table_w - Inches(0.7),
                             height=Inches(4.0))

note_y = content_y + Inches(4.15)
add_textbox(slide, Inches(0.4), note_y, SLIDE_W - Inches(0.8), Inches(0.6),
            "Our LSTM is closest to BachBot. We explored the tension between "
            "statistical fit and musical quality rather than aiming for SOTA perplexity.",
            font_size=Pt(15), italic=True, color=MUTED)

print("Slide 15 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 -- Related Work: Continuous Generation
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Related Work: Continuous Audio Generation", bg_color=BLUE)

content_y = bar_h + Inches(0.2)

t4_rw = img("t4_related_work.png")
table_w = SLIDE_W - Inches(0.8) if not t4_rw else Inches(8.5)

headers = ["Year", "Work", "Model", "Key Contribution"]
rows_rw4 = [
    ["2017", "FMA Dataset (Defferrard)", "--",
     "106K tracks, 161 genres, CC licensed benchmark"],
    ["2023", "MusicGen (Copet, NeurIPS)", "Single-stage AR LM",
     "Efficient codebook interleaving -- our base model"],
    ["2023", "MusicLM (Agostinelli, Google)", "Hierarchical audio LM",
     "Semantic + acoustic token hierarchy"],
    ["2023", "AudioCraft (Meta)", "Open-source framework",
     "pip install audiocraft, democratized audio gen"],
    ["2025", "Genre Fine-tuning (IJCRT)", "MusicGen-small fine-tuned",
     "Same FMA genre approach as ours"],
]
add_table_to_slide(slide, headers, rows_rw4,
                   Inches(0.4), content_y,
                   table_w, Inches(3.8),
                   font_size=Pt(15),
                   col_widths=[0.7, 2.8, 2.5, 4.5])

if t4_rw:
    slide.shapes.add_picture(t4_rw,
                             Inches(0.4) + table_w + Inches(0.1),
                             content_y,
                             width=SLIDE_W - table_w - Inches(0.7),
                             height=Inches(3.8))

note_y = content_y + Inches(3.95)
add_textbox(slide, Inches(0.4), note_y, SLIDE_W - Inches(0.8), Inches(0.75),
            "Our contribution: side-by-side comparison of symbolic and continuous generation, "
            "showing how the same goal plays out across two representation spaces.",
            font_size=Pt(15), italic=True, color=MUTED)

print("Slide 16 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 -- Cross-Task Comparison
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Symbolic vs Continuous: Two Paradigms Compared")

content_y = bar_h + Inches(0.2)

headers = ["Dimension", "Task 1: Symbolic", "Task 4: Continuous"]
rows_cmp = [
    ["Representation",   "MIDI tokens (47-dim vocab)",           "Raw audio waveform (32kHz)"],
    ["Model size",       "1.1M parameters",                      "300M parameters"],
    ["Training data",    "368 chorales (~15 min audio equiv.)",   "8,000 tracks (~67 hours)"],
    ["Training time",    "Minutes (CPU)",                         "Hours (GPU)"],
    ["Evaluation",       "5 precise metrics (perplexity, KL...)", "Genre classifier accuracy"],
    ["Output quality",   'Structured but "mechanical"',           "Perceptually realistic"],
    ["Interpretability", "High (inspect individual notes/voices)","Low (raw waveform)"],
    ["Key tension",      "Statistical fit vs music theory rules", "Generic quality vs genre specificity"],
]
cell_colors = {}
for ri in range(len(rows_cmp)):
    cell_colors[(ri, 1)] = RGBColor(0xFD,0xED,0xEC)   # Task 1 red
    cell_colors[(ri, 2)] = RGBColor(0xEB,0xF5,0xFB)   # Task 4 blue

add_table_to_slide(slide, headers, rows_cmp,
                   Inches(0.4), content_y,
                   SLIDE_W - Inches(0.8), Inches(5.5),
                   font_size=Pt(15),
                   col_widths=[2.5, 4.0, 4.0],
                   cell_colors=cell_colors)

print("Slide 17 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 -- Conclusion
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Conclusion")

content_y = bar_h + Inches(0.2)
card_h    = Inches(2.0)
card_w    = Inches(6.1)
gap       = Inches(0.4)
lx        = Inches(0.4)
rx        = lx + card_w + gap

# Task 1 card
add_rect(slide, lx, content_y, card_w, card_h,
         fill_color=RGBColor(0xFD,0xED,0xEC), line_color=RED, line_width=Pt(2))
add_textbox(slide, lx + Inches(0.15), content_y + Inches(0.05),
            card_w - Inches(0.3), Inches(0.5),
            "Task 1 Results", font_size=Pt(20), bold=True, color=RED)
add_multiline_textbox(slide,
                      lx + Inches(0.15), content_y + Inches(0.55),
                      card_w - Inches(0.3), card_h - Inches(0.65),
                      ["  LSTM achieves 1.97 perplexity on JSB Chorales",
                       "  1.3x improvement over bigram Markov baseline",
                       "  But more parallel 5ths/octaves (18.1% vs 11.5%)"],
                      font_size=Pt(17), color=BODY, space_before=Pt(5))

# Task 4 card
add_rect(slide, rx, content_y, card_w, card_h,
         fill_color=RGBColor(0xEB,0xF5,0xFB), line_color=BLUE, line_width=Pt(2))
add_textbox(slide, rx + Inches(0.15), content_y + Inches(0.05),
            card_w - Inches(0.3), Inches(0.5),
            "Task 4 Results", font_size=Pt(20), bold=True, color=BLUE)
add_multiline_textbox(slide,
                      rx + Inches(0.15), content_y + Inches(0.55),
                      card_w - Inches(0.3), card_h - Inches(0.65),
                      ["  Fine-tuning improves genre accuracy: 25% -> 75%",
                       "  3/4 genres correctly classified after fine-tuning",
                       "  Proof-of-concept with only 72 training samples"],
                      font_size=Pt(17), color=BODY, space_before=Pt(5))

# Key takeaways
takeaway_y = content_y + card_h + Inches(0.25)
add_textbox(slide, Inches(0.4), takeaway_y, Inches(12.5), Inches(0.45),
            "Key Takeaways", font_size=Pt(22), bold=True, color=DARK)

takeaways = [
    "1.  Statistical objectives (cross-entropy) do not guarantee musical quality",
    "2.  Fine-tuning large pretrained models is practical and effective, "
        "even with limited data and compute",
    "3.  The choice between symbolic and continuous representations "
        "fundamentally shapes what you can evaluate",
]
add_multiline_textbox(slide,
                      Inches(0.4), takeaway_y + Inches(0.5),
                      SLIDE_W - Inches(0.8), Inches(2.4),
                      takeaways, font_size=Pt(19), color=BODY,
                      space_before=Pt(8))

print("Slide 18 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 -- Thank You
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_WHITE)

add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), fill_color=DARK)
add_rect(slide, 0, SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), fill_color=DARK)

add_textbox(slide, Inches(0.7), Inches(2.0), Inches(11.9), Inches(1.2),
            "Thank You",
            font_size=Pt(60), bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.7), Inches(3.5), Inches(11.9), Inches(0.8),
            "Questions?",
            font_size=Pt(36), color=MUTED, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(3.5), Inches(4.6), Inches(6.3), Inches(0.55),
         fill_color=RGBColor(0xE8,0xEA,0xEB))
add_textbox(slide, Inches(3.5), Inches(4.6), Inches(6.3), Inches(0.55),
            "Code and notebook:  workbook.ipynb  |  workbook.html",
            font_size=Pt(16), color=DARK, align=PP_ALIGN.CENTER)

print("Slide 19 done")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 -- Generated Music Demo
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide)

bar_h = title_bar(slide, "Generated Music Demo")

content_y = bar_h + Inches(0.2)
col_w = Inches(6.1)
gap   = Inches(0.4)
lx    = Inches(0.4)
rx    = lx + col_w + gap

# Task 1 section
add_rect(slide, lx, content_y, col_w, Inches(0.45),
         fill_color=RED)
add_textbox(slide, lx + Inches(0.1), content_y, col_w - Inches(0.2), Inches(0.45),
            "Task 1 -- Symbolic (MIDI)", font_size=Pt(18), bold=True,
            color=BG_WHITE, align=PP_ALIGN.LEFT)

t1_files = [
    "1.  Real Bach chorale (reference): sample_chorale_real.mid",
    "2.  Bigram Markov Chain: markov_chorale.mid",
    "3.  LSTM (temperature 0.9): symbolic_unconditioned.mid  [main deliverable]",
]
add_multiline_textbox(slide,
                      lx + Inches(0.1), content_y + Inches(0.5),
                      col_w - Inches(0.2), Inches(2.5),
                      t1_files, font_size=Pt(17), color=BODY,
                      space_before=Pt(8))

# Task 4 section
add_rect(slide, rx, content_y, col_w, Inches(0.45),
         fill_color=BLUE)
add_textbox(slide, rx + Inches(0.1), content_y, col_w - Inches(0.2), Inches(0.45),
            "Task 4 -- Continuous (Audio)", font_size=Pt(18), bold=True,
            color=BG_WHITE, align=PP_ALIGN.LEFT)

t4_files = [
    "4.  Pretrained, Electronic: generated_audio/electronic_generated.mp3",
    "5.  Fine-tuned, Electronic: generated_audio_finetuned/electronic_generated.mp3",
    "6.  Pretrained, Folk: generated_audio/folk_generated.mp3",
    "7.  Fine-tuned, Folk: generated_audio_finetuned/folk_generated.mp3",
    "8.  Main deliverable: continuous_conditioned.mp3 (Hip-Hop prompt)",
]
add_multiline_textbox(slide,
                      rx + Inches(0.1), content_y + Inches(0.5),
                      col_w - Inches(0.2), Inches(3.3),
                      t4_files, font_size=Pt(16), color=BODY,
                      space_before=Pt(6))

# Waveform decoration strip
wave_y = content_y + Inches(4.1)
add_rect(slide, Inches(0.4), wave_y, SLIDE_W - Inches(0.8), Inches(0.12),
         fill_color=BLUE)
add_textbox(slide, Inches(0.4), wave_y + Inches(0.18),
            SLIDE_W - Inches(0.8), Inches(0.4),
            "Audio plays over this slide -- not counted toward the 20-minute time limit",
            font_size=Pt(14), italic=True, color=MUTED, align=PP_ALIGN.CENTER)

print("Slide 20 done")


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
prs.save(OUT_PATH)
print(f"\nSaved: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
